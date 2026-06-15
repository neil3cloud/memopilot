"""Investigation mode services for evidence-aware analysis."""

from __future__ import annotations

import csv
import importlib.util
import json
import re
import uuid
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

import aiosqlite
import openpyxl
import pdfplumber
from docx import Document
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pydantic import BaseModel, Field

from .config import Config
from .db import DatabaseManager
from .evidence_classifier import EvidenceSourceClassifier
from .plan_service import PlanModeService, PlanResult, PlanStep
from .security_policy import CredentialRedactor
from .workspace_roots import WorkspaceRootsService

MAX_EVIDENCE_FILE_SIZE_BYTES = 10 * 1024 * 1024


class AttachedEvidence(BaseModel):
    evidence_id: str
    source_type: str
    trust_level: int
    extraction_method: str
    extraction_status: str
    findings: list[str] = Field(default_factory=list)
    redacted_values: int
    source_path: str | None
    investigation_session_id: str | None = None


class EvidenceBoardEntry(BaseModel):
    evidence_id: str
    source_type: str
    source_path: str | None
    source_url: str | None
    trust_level: int
    extraction_method: str
    extraction_status: str
    redacted_values: int
    findings: list[str] = Field(default_factory=list)
    investigation_session_id: str | None = None


class EvidenceColumnsPreview(BaseModel):
    source_type: str
    columns: list[str]
    suggested_mapping: dict[str, str]
    requires_confirmation: bool


class InvestigationResult(BaseModel):
    context_pack: str
    context_pack_path: str
    impacted_files: list[str]
    related_tests: list[str]
    missing_test_coverage: list[str]
    evidence_count: int


class InvestigationSessionCreate(BaseModel):
    title: str
    description: str = ""
    mode: str = "investigation"


class InvestigationSession(BaseModel):
    id: str
    title: str
    description: str | None = None
    mode: str
    status: str
    workspace_root: str
    created_at: str
    updated_at: str
    evidence_count: int = 0
    evidence: list[EvidenceBoardEntry] = Field(default_factory=list)


class RemovedEvidence(BaseModel):
    evidence_id: str
    removed: bool = True


@dataclass(frozen=True)
class InvestigationFindingsSummary:
    findings: list[str]
    root_cause: list[str]
    impacted_files: list[str]
    impacted_callers: list[str]
    acceptance_criteria: list[str]
    task_run_id: str | None = None


class InvestigationService:
    """Implements evidence attachment and investigation context pack generation."""

    def __init__(self, *, config: Config, db: DatabaseManager) -> None:
        self._config = config
        self._db = db
        self._redactor = CredentialRedactor()
        self._classifier = EvidenceSourceClassifier()

    async def preview_columns(
        self,
        *,
        evidence_path: str,
        workspace_root: str | None = None,
    ) -> EvidenceColumnsPreview:
        resolved_path = await self._resolve_evidence_path(
            evidence_path, workspace_root=workspace_root
        )
        source_type = self._classifier.classify(
            evidence_path=resolved_path,
            source_url=None,
        ).source_type

        if source_type == "csv_data":
            headers, _rows = self._extract_csv_rows(resolved_path)
            mapping = self._suggest_column_mapping(headers)
            return EvidenceColumnsPreview(
                source_type=source_type,
                columns=headers,
                suggested_mapping=mapping,
                requires_confirmation=True,
            )
        if source_type == "spreadsheet":
            headers, _rows, _sheet = self._extract_excel_rows(resolved_path)
            mapping = self._suggest_column_mapping(headers)
            return EvidenceColumnsPreview(
                source_type=source_type,
                columns=headers,
                suggested_mapping=mapping,
                requires_confirmation=True,
            )

        return EvidenceColumnsPreview(
            source_type=source_type,
            columns=[],
            suggested_mapping={},
            requires_confirmation=False,
        )

    async def start_session(
        self,
        *,
        title: str,
        description: str = "",
        mode: str = "investigation",
        workspace_root: str | None = None,
    ) -> InvestigationSession:
        clean_title = title.strip()
        if not clean_title:
            raise ValueError("Investigation title is required")
        clean_mode = mode.strip() if mode and mode.strip() else "investigation"
        workspace_service = WorkspaceRootsService(config=self._config, db=self._db)
        workspace_root = str(await workspace_service.resolve_workspace_root(workspace_root))
        session_id = uuid.uuid4().hex
        conn = await self._db.connect()
        await conn.execute(
            """
            INSERT INTO investigation_sessions
            (
                id,
                title,
                description,
                mode,
                status,
                workspace_root,
                created_at,
                updated_at
            )
            VALUES (?, ?, ?, ?, 'open', ?, datetime('now'), datetime('now'))
            """,
            (session_id, clean_title, description.strip(), clean_mode, workspace_root),
        )
        await conn.commit()
        return await self.get_session(session_id=session_id, conn=conn)

    async def get_session(
        self,
        *,
        session_id: str,
        conn: aiosqlite.Connection | None = None,
    ) -> InvestigationSession:
        connection = conn or await self._db.connect()
        row = await self._fetch_session_row(connection, session_id=session_id)
        if row is None:
            raise ValueError(f"Investigation session not found: {session_id}")
        evidence = await self.list_evidence_board(
            investigation_session_id=session_id,
            task_run_id=None,
            workspace_root=str(row["workspace_root"]),
            conn=connection,
        )
        return self._build_session(row, evidence)

    async def attach_evidence(
        self,
        *,
        evidence_path: str | None,
        source_url: str | None,
        task_run_id: str | None,
        investigation_session_id: str | None = None,
        column_mapping: dict[str, str] | None = None,
        workspace_root: str | None = None,
    ) -> AttachedEvidence:
        conn = await self._db.connect()
        session_root: str | None = None
        if investigation_session_id is not None:
            session_row = await self._require_session(conn, session_id=investigation_session_id)
            session_root = str(session_row["workspace_root"])
        effective_workspace_root = workspace_root or session_root
        resolved_path = (
            await self._resolve_evidence_path(
                evidence_path, workspace_root=effective_workspace_root
            )
            if evidence_path
            else None
        )
        classification = self._classifier.classify(
            evidence_path=resolved_path,
            source_url=source_url,
            content_preview=self._read_content_preview(resolved_path),
        )
        source_type = classification.source_type
        extraction_method = classification.extraction_method
        trust_level = classification.trust_level

        findings, extraction_status, ocr_required = self._extract_findings(
            source_type=source_type,
            evidence_path=resolved_path,
            column_mapping=column_mapping or {},
        )
        if ocr_required and source_type == "pdf_doc":
            trust_level = 4

        redacted_findings, redacted_values = self._redact_findings(findings)

        evidence_id = uuid.uuid4().hex
        await conn.execute(
            """
            INSERT INTO evidence_sources
            (
                id,
                task_run_id,
                investigation_session_id,
                source_type,
                source_path,
                source_url,
                trust_level,
                extraction_method,
                extracted_findings_json,
                approved,
                workspace_root
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, 0, ?)
            """,
            (
                evidence_id,
                task_run_id,
                investigation_session_id,
                source_type,
                str(resolved_path) if resolved_path else None,
                source_url,
                trust_level,
                extraction_method,
                json.dumps(
                    {
                        "findings": redacted_findings,
                        "extraction_status": extraction_status,
                        "redacted_values": redacted_values,
                    }
                ),
                effective_workspace_root,
            ),
        )
        await conn.commit()

        return AttachedEvidence(
            evidence_id=evidence_id,
            source_type=source_type,
            trust_level=trust_level,
            extraction_method=extraction_method,
            extraction_status=extraction_status,
            findings=redacted_findings,
            redacted_values=redacted_values,
            source_path=str(resolved_path) if resolved_path else None,
            investigation_session_id=investigation_session_id,
        )

    async def remove_evidence(
        self,
        *,
        session_id: str,
        evidence_id: str,
    ) -> RemovedEvidence:
        conn = await self._db.connect()
        await self._require_session(conn, session_id=session_id)
        cursor = await conn.execute(
            """
            SELECT id
            FROM evidence_sources
            WHERE id = ? AND investigation_session_id = ?
            LIMIT 1
            """,
            (evidence_id, session_id),
        )
        row = await cursor.fetchone()
        if row is None:
            raise ValueError(f"Evidence not found for investigation session: {evidence_id}")
        await conn.execute(
            "DELETE FROM evidence_sources WHERE id = ? AND investigation_session_id = ?",
            (evidence_id, session_id),
        )
        await conn.commit()
        return RemovedEvidence(evidence_id=evidence_id)

    async def transition_to_patch(self, *, session_id: str) -> InvestigationSession:
        conn = await self._db.connect()
        await self._require_session(conn, session_id=session_id)
        await conn.execute(
            """
            UPDATE investigation_sessions
            SET status = 'patch_generated', updated_at = datetime('now')
            WHERE id = ?
            """,
            (session_id,),
        )
        await conn.commit()
        return await self.get_session(session_id=session_id, conn=conn)

    async def list_evidence_board(
        self,
        *,
        task_run_id: str | None,
        investigation_session_id: str | None = None,
        workspace_root: str | None = None,
        conn: aiosqlite.Connection | None = None,
    ) -> list[EvidenceBoardEntry]:
        connection = conn or await self._db.connect()
        params: list[str] = []
        where_clauses: list[str] = []
        if investigation_session_id:
            where_clauses.append("investigation_session_id = ?")
            params.append(investigation_session_id)
        if task_run_id:
            where_clauses.append("task_run_id = ?")
            params.append(task_run_id)
        if workspace_root:
            where_clauses.append("COALESCE(workspace_root, ?) = ?")
            params.extend([workspace_root, workspace_root])

        query = """
            SELECT
                id,
                investigation_session_id,
                source_type,
                source_path,
                source_url,
                trust_level,
                extraction_method,
                extracted_findings_json
            FROM evidence_sources
        """
        if where_clauses:
            query += " WHERE " + " AND ".join(where_clauses)
        query += " ORDER BY created_at DESC"
        if not where_clauses:
            query += " LIMIT 100"

        cursor = await connection.execute(query, tuple(params))
        rows = await cursor.fetchall()
        return [
            EvidenceBoardEntry(
                evidence_id=row["id"],
                source_type=row["source_type"],
                source_path=row["source_path"],
                source_url=row["source_url"],
                trust_level=int(row["trust_level"]),
                extraction_method=row["extraction_method"] or "unknown",
                extraction_status=self._parse_extraction_status(
                    row["source_type"],
                    row["extracted_findings_json"],
                ),
                redacted_values=self._parse_redacted_values(row["extracted_findings_json"]),
                findings=self._parse_findings(row["extracted_findings_json"]),
                investigation_session_id=row["investigation_session_id"],
            )
            for row in rows
        ]

    async def run_investigation(
        self,
        *,
        title: str,
        description: str,
        acceptance_criteria: list[str],
        task_run_id: str | None,
        workspace_root: str | None = None,
    ) -> InvestigationResult:
        evidence_entries = await self.list_evidence_board(
            task_run_id=task_run_id,
            workspace_root=workspace_root,
        )
        combined_findings = [finding for entry in evidence_entries for finding in entry.findings]
        impacted_files = await self._discover_impacted_files(combined_findings)
        related_tests = await self._discover_related_tests(impacted_files)
        missing_coverage = self._detect_missing_coverage(acceptance_criteria, related_tests)
        active_rules = await self._active_rules()

        context_pack = self._build_context_pack(
            title=title,
            description=description,
            acceptance_criteria=acceptance_criteria,
            evidence_entries=evidence_entries,
            impacted_files=impacted_files,
            related_tests=related_tests,
            missing_coverage=missing_coverage,
            active_rules=active_rules,
        )
        context_pack_path = self._write_context_pack(context_pack)
        return InvestigationResult(
            context_pack=context_pack,
            context_pack_path=context_pack_path,
            impacted_files=impacted_files,
            related_tests=related_tests,
            missing_test_coverage=missing_coverage,
            evidence_count=len(evidence_entries),
        )

    async def generate_plan_from_findings(
        self,
        *,
        investigation_session_id: str,
        workspace_root: str | None = None,
    ) -> PlanResult:
        conn = await self._db.connect()
        session_row = await self._require_session(conn, session_id=investigation_session_id)
        normalized_root = self._normalize_workspace_root(
            workspace_root or session_row["workspace_root"]
        )
        summary = await self._summarize_findings(
            conn=conn,
            investigation_session_id=investigation_session_id,
            workspace_root=normalized_root,
        )
        steps = self._build_plan_steps(summary)
        if not steps:
            raise ValueError("No actionable investigation findings available to build a plan")

        plan_service = PlanModeService(config=self._config, db=self._db)
        plan = await plan_service.store_plan(
            title=f"Plan from investigation: {session_row['title']}",
            steps=steps,
            task_description=self._build_plan_task_description(session_row, summary),
            workspace_root=normalized_root,
            task_run_id=summary.task_run_id,
        )
        await self.store_investigation_memory(
            investigation_session_id=investigation_session_id,
            workspace_root=normalized_root,
            summary=summary,
        )
        return plan

    async def store_investigation_memory(
        self,
        *,
        investigation_session_id: str,
        workspace_root: str | None = None,
        summary: InvestigationFindingsSummary | None = None,
    ) -> list[str]:
        conn = await self._db.connect()
        session_row = await self._require_session(conn, session_id=investigation_session_id)
        normalized_root = self._normalize_workspace_root(
            workspace_root or session_row["workspace_root"]
        )
        findings_summary = summary or await self._summarize_findings(
            conn=conn,
            investigation_session_id=investigation_session_id,
            workspace_root=normalized_root,
        )

        memory_ids: list[str] = []
        root_cause_lines = findings_summary.root_cause or findings_summary.findings[:1]
        if root_cause_lines:
            memory_ids.append(
                await self._store_confirmed_memory_item(
                    conn=conn,
                    investigation_session_id=investigation_session_id,
                    workspace_root=normalized_root,
                    finding_kind="root_cause",
                    title=f"Root cause for investigation: {session_row['title']}",
                    body=self._bullet_list(root_cause_lines),
                    memory_class="fact",
                    trust_level=4,
                )
            )

        impacted_lines: list[str] = []
        if findings_summary.impacted_files:
            impacted_lines.append("Impacted files:")
            impacted_lines.extend(f"- {path}" for path in findings_summary.impacted_files)
        if findings_summary.impacted_callers:
            impacted_lines.append("Impacted callers:")
            impacted_lines.extend(f"- {caller}" for caller in findings_summary.impacted_callers)
        if impacted_lines:
            memory_ids.append(
                await self._store_confirmed_memory_item(
                    conn=conn,
                    investigation_session_id=investigation_session_id,
                    workspace_root=normalized_root,
                    finding_kind="impacted_scope",
                    title=f"Impacted scope for investigation: {session_row['title']}",
                    body="\n".join(impacted_lines),
                    memory_class="fact",
                    trust_level=4,
                )
            )

        if findings_summary.acceptance_criteria:
            memory_ids.append(
                await self._store_confirmed_memory_item(
                    conn=conn,
                    investigation_session_id=investigation_session_id,
                    workspace_root=normalized_root,
                    finding_kind="acceptance_criteria",
                    title=f"Acceptance criteria for investigation: {session_row['title']}",
                    body=self._bullet_list(findings_summary.acceptance_criteria),
                    memory_class="instruction",
                    trust_level=3,
                )
            )

        await conn.commit()
        return memory_ids

    async def _fetch_session_row(
        self,
        conn: aiosqlite.Connection,
        *,
        session_id: str,
    ) -> aiosqlite.Row | None:
        cursor = await conn.execute(
            """
            SELECT
                id,
                title,
                description,
                mode,
                status,
                workspace_root,
                created_at,
                updated_at
            FROM investigation_sessions
            WHERE id = ?
            LIMIT 1
            """,
            (session_id,),
        )
        return await cursor.fetchone()

    async def _require_session(
        self,
        conn: aiosqlite.Connection,
        *,
        session_id: str,
    ) -> aiosqlite.Row:
        row = await self._fetch_session_row(conn, session_id=session_id)
        if row is None:
            raise ValueError(f"Investigation session not found: {session_id}")
        return row

    def _build_session(
        self,
        row: aiosqlite.Row,
        evidence: list[EvidenceBoardEntry],
    ) -> InvestigationSession:
        return InvestigationSession(
            id=row["id"],
            title=row["title"],
            description=row["description"],
            mode=row["mode"],
            status=row["status"],
            workspace_root=row["workspace_root"],
            created_at=row["created_at"],
            updated_at=row["updated_at"],
            evidence_count=len(evidence),
            evidence=evidence,
        )

    async def _resolve_evidence_path(
        self,
        evidence_path: str,
        *,
        workspace_root: str | None = None,
    ) -> Path:
        candidate = Path(evidence_path)
        workspace_service = WorkspaceRootsService(config=self._config, db=self._db)
        scoped_root = await workspace_service.resolve_workspace_root(workspace_root)
        allowed_roots = (
            [scoped_root] if workspace_root else await workspace_service.allowed_workspace_paths()
        )
        if not candidate.is_absolute():
            if self._contains_parent_reference(candidate):
                raise ValueError("Evidence path must not traverse parent directories")
            candidate = scoped_root / candidate
        if not any(self._is_within_workspace(candidate, root) for root in allowed_roots):
            raise ValueError("Evidence path must be within a configured workspace root")
        if candidate.is_symlink():
            raise ValueError("Symlinked evidence paths are not allowed")
        resolved = candidate.resolve()
        if not any(self._is_within_workspace(resolved, root) for root in allowed_roots):
            raise ValueError("Evidence path must be within a configured workspace root")
        if not resolved.exists():
            raise ValueError(f"Evidence path not found: {resolved}")
        return resolved

    def _extract_findings(
        self,
        *,
        source_type: str,
        evidence_path: Path | None,
        column_mapping: dict[str, str],
    ) -> tuple[list[str], str, bool]:
        if source_type == "external_work_item":
            return (
                ["External work item attached. Fetch details via MCP in future iterations."],
                "ok",
                False,
            )
        if evidence_path is None:
            return ["No file evidence provided."], "ok", False
        # Guard against excessively large files (10 MB limit)
        try:
            file_size = evidence_path.stat().st_size
        except OSError as exc:
            raise ValueError(f"Cannot access evidence file: {exc}") from exc
        if file_size > MAX_EVIDENCE_FILE_SIZE_BYTES:
            raise ValueError(
                f"Evidence file too large ({file_size} bytes, max {MAX_EVIDENCE_FILE_SIZE_BYTES})"
            )
        if source_type in {"image", "screenshot"}:
            return self._extract_image_findings(evidence_path, source_type)
        if source_type == "csv_data":
            headers, rows = self._extract_csv_rows(evidence_path)
            mapping = self._normalize_column_mapping(column_mapping, headers)
            return (
                self._build_tabular_findings(
                    source_label=f"CSV:{evidence_path.name}",
                    headers=headers,
                    rows=rows,
                    mapping=mapping,
                ),
                "ok",
                False,
            )
        if source_type == "spreadsheet":
            headers, rows, sheet_name = self._extract_excel_rows(evidence_path)
            mapping = self._normalize_column_mapping(column_mapping, headers)
            return (
                self._build_tabular_findings(
                    source_label=f"Excel:{evidence_path.name}:{sheet_name}",
                    headers=headers,
                    rows=rows,
                    mapping=mapping,
                ),
                "ok",
                False,
            )
        if source_type == "pdf_doc":
            return self._extract_pdf_findings(evidence_path)
        if source_type == "word_doc":
            return self._extract_word_findings(evidence_path), "ok", False
        if source_type == "powerpoint_doc":
            return self._extract_powerpoint_findings(evidence_path), "ok", False

        try:
            text = evidence_path.read_text(encoding="utf-8", errors="replace")
        except OSError as exc:
            raise ValueError(f"Cannot read evidence file: {exc}") from exc
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return lines[:30], "ok", False

    def _read_content_preview(self, evidence_path: Path | None) -> str | None:
        if evidence_path is None:
            return None
        try:
            with evidence_path.open("rb") as handle:
                return handle.read(8192).decode("utf-8", errors="replace")
        except OSError:
            return None

    def _extract_image_findings(
        self,
        evidence_path: Path,
        source_type: str,
    ) -> tuple[list[str], str, bool]:
        findings: list[str] = []
        try:
            with Image.open(evidence_path) as image:
                findings.append(f"{source_type.title()} evidence: {evidence_path.name}")
                image_format = image.format or evidence_path.suffix.lstrip(".").upper()
                findings.append(
                    f"Image analysis: {image.width}x{image.height}, "
                    f"mode={image.mode}, format={image_format}"
                )

                if source_type == "screenshot":
                    findings.append(
                        "Screenshot heuristic: filename suggests a captured UI or screen state."
                    )

                metadata = self._image_metadata(image)
                if metadata:
                    findings.append(f"Image metadata keys: {', '.join(metadata)}")

                ocr_lines = self._extract_image_text(image)
                if ocr_lines:
                    findings.extend(ocr_lines)
                    return findings, "ok", False

                findings.append("OCR unavailable; retained metadata-based image analysis.")
                return findings, "metadata_only", False
        except UnidentifiedImageError:
            return (
                [f"Image file could not be decoded: {evidence_path.name}"],
                "unreadable_image",
                False,
            )

    def _image_metadata(self, image: Image.Image) -> list[str]:
        keys = [str(key) for key in image.getexif().keys()][:8]
        if not keys:
            return []
        return sorted(keys)

    def _extract_image_text(self, image: Image.Image) -> list[str]:
        if importlib.util.find_spec("pytesseract") is None:
            return []

        import pytesseract

        try:
            raw_text = pytesseract.image_to_string(image).strip()
        except Exception:
            return []

        lines = [line.strip() for line in raw_text.splitlines() if line.strip()]
        return [f"OCR: {line}" for line in lines[:20]]

    def _extract_csv_rows(self, evidence_path: Path) -> tuple[list[str], list[dict[str, str]]]:
        with evidence_path.open("r", encoding="utf-8", errors="replace", newline="") as file_obj:
            reader = csv.DictReader(file_obj)
            headers = [header for header in (reader.fieldnames or []) if header]
            rows: list[dict[str, str]] = []
            for index, row in enumerate(reader):
                if index >= 50:
                    break
                rows.append({str(key): str(value) for key, value in row.items() if key and value})
        return headers, rows

    def _extract_excel_rows(
        self,
        evidence_path: Path,
    ) -> tuple[list[str], list[dict[str, str]], str]:
        workbook = openpyxl.load_workbook(evidence_path, read_only=True, data_only=True)
        worksheet = workbook.active
        rows_iter = worksheet.iter_rows(values_only=True)
        first_row = next(rows_iter, None)
        headers = [str(value).strip() for value in (first_row or []) if value is not None]

        rows: list[dict[str, str]] = []
        for index, row in enumerate(rows_iter):
            if index >= 50:
                break
            entry: dict[str, str] = {}
            for column_index, value in enumerate(row):
                if column_index >= len(headers):
                    continue
                if value is None:
                    continue
                entry[headers[column_index]] = str(value).strip()
            if entry:
                rows.append(entry)
        workbook.close()
        return headers, rows, worksheet.title

    def _extract_pdf_findings(self, evidence_path: Path) -> tuple[list[str], str, bool]:
        findings: list[str] = []
        saw_text = False
        with pdfplumber.open(evidence_path) as pdf_doc:
            for index, page in enumerate(pdf_doc.pages):
                if index >= 20:
                    break
                page_text = (page.extract_text() or "").strip()
                if not page_text:
                    continue
                saw_text = True
                for line in page_text.splitlines():
                    stripped = line.strip()
                    if stripped:
                        findings.append(stripped)
                if len(findings) >= 60:
                    break
        if not saw_text:
            return (
                ["Scanned PDF detected. OCR is required before reliable extraction."],
                "requires_ocr",
                True,
            )
        return findings[:60], "ok", False

    def _extract_word_findings(self, evidence_path: Path) -> list[str]:
        document = Document(evidence_path)
        findings = [
            paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()
        ]
        if not findings:
            return [f"Word document contains no readable paragraph text: {evidence_path.name}"]
        return findings[:80]

    def _extract_powerpoint_findings(self, evidence_path: Path) -> list[str]:
        presentation = Presentation(evidence_path)
        findings: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if not text:
                    continue
                for line in text.splitlines():
                    stripped = line.strip()
                    if stripped:
                        findings.append(f"Slide {slide_index}: {stripped}")
                if len(findings) >= 80:
                    break
            if len(findings) >= 80:
                break
        if not findings:
            return [f"PowerPoint contains no readable slide text: {evidence_path.name}"]
        return findings

    def _build_tabular_findings(
        self,
        *,
        source_label: str,
        headers: list[str],
        rows: list[dict[str, str]],
        mapping: dict[str, str],
    ) -> list[str]:
        findings: list[str] = []
        findings.append(f"Tabular source: {source_label}")
        if headers:
            findings.append(f"Columns: {', '.join(headers)}")

        for column in headers:
            values = [row[column] for row in rows if column in row and row[column]]
            sample_type = self._infer_value_type(values)
            findings.append(
                f"Data dictionary: {column} (type={sample_type}, sample_count={len(values[:20])})"
            )

        if not mapping and headers:
            mapping = self._suggest_column_mapping(headers)

        title_col = mapping.get("title")
        steps_col = mapping.get("steps")
        expected_col = mapping.get("expected")
        if title_col:
            for row in rows[:15]:
                title = row.get(title_col, "").strip()
                steps = row.get(steps_col, "").strip() if steps_col else ""
                expected = row.get(expected_col, "").strip() if expected_col else ""
                if not title:
                    continue
                findings.append(
                    "Test case: "
                    f"title={title}; steps={steps or 'n/a'}; expected={expected or 'n/a'}"
                )

        if not rows:
            findings.append("No non-empty tabular rows were parsed.")
        return findings

    def _suggest_column_mapping(self, headers: list[str]) -> dict[str, str]:
        lowered = {header.lower(): header for header in headers}

        def pick(candidates: tuple[str, ...]) -> str | None:
            for candidate in candidates:
                for lower, original in lowered.items():
                    if candidate in lower:
                        return original
            return None

        mapping: dict[str, str] = {}
        title = pick(("title", "name", "summary", "test case"))
        steps = pick(("steps", "procedure", "action", "input"))
        expected = pick(("expected", "result", "assert", "outcome"))
        if title:
            mapping["title"] = title
        if steps:
            mapping["steps"] = steps
        if expected:
            mapping["expected"] = expected
        return mapping

    def _normalize_column_mapping(
        self,
        mapping: dict[str, str],
        headers: list[str],
    ) -> dict[str, str]:
        if not mapping:
            return {}
        index = {header.lower(): header for header in headers}
        normalized: dict[str, str] = {}
        for role, column in mapping.items():
            resolved = index.get(column.lower())
            if resolved:
                normalized[role] = resolved
        return normalized

    def _infer_value_type(self, values: list[str]) -> str:
        sample = [value for value in values[:20] if value]
        if not sample:
            return "empty"
        if all(re.fullmatch(r"-?\d+(\.\d+)?", value) for value in sample):
            return "number"
        if all(re.fullmatch(r"\d{4}-\d{2}-\d{2}", value) for value in sample):
            return "date"
        if all(value.lower() in {"true", "false", "yes", "no"} for value in sample):
            return "boolean"
        return "text"

    def _redact_findings(self, findings: list[str]) -> tuple[list[str], int]:
        total = 0
        redacted: list[str] = []
        for finding in findings:
            result = self._redactor.redact(finding)
            redacted.append(result.redacted_text)
            total += result.redacted_count
        return redacted, total

    async def _discover_impacted_files(self, findings: list[str]) -> list[str]:
        if not findings:
            return []
        tokens = self._important_tokens(findings)
        if not tokens:
            return []

        conn = await self._db.connect()
        file_cursor = await conn.execute("SELECT file_path FROM file_index LIMIT 10000")
        symbol_cursor = await conn.execute(
            "SELECT DISTINCT file_path, name FROM symbols LIMIT 10000"
        )
        file_rows = await file_cursor.fetchall()
        symbol_rows = await symbol_cursor.fetchall()

        matched: set[str] = set()
        for row in file_rows:
            file_path = str(row["file_path"])
            lowered = file_path.lower()
            if any(token in lowered for token in tokens):
                matched.add(file_path)

        for row in symbol_rows:
            file_path = str(row["file_path"])
            symbol_name = str(row["name"]).lower()
            if any(token in symbol_name for token in tokens):
                matched.add(file_path)

        return sorted(matched)[:100]

    async def _discover_related_tests(self, impacted_files: list[str]) -> list[str]:
        conn = await self._db.connect()
        cursor = await conn.execute(
            "SELECT file_path FROM file_index WHERE file_path LIKE '%test%'"
        )
        rows = await cursor.fetchall()
        candidates = [str(row["file_path"]) for row in rows]
        impacted_tokens = {
            token
            for path in impacted_files
            for token in re.split(r"[^a-zA-Z0-9_]+", path.lower())
            if len(token) >= 4
        }
        related = [
            test_file
            for test_file in candidates
            if impacted_tokens.intersection(re.split(r"[^a-zA-Z0-9_]+", test_file.lower()))
        ]
        return sorted(set(related))[:100]

    def _detect_missing_coverage(
        self,
        acceptance_criteria: list[str],
        related_tests: list[str],
    ) -> list[str]:
        if not acceptance_criteria:
            return []
        if not related_tests:
            return acceptance_criteria
        joined_tests = " ".join(related_tests).lower()
        missing: list[str] = []
        generic_tokens = {
            "should",
            "must",
            "when",
            "then",
            "with",
            "from",
            "into",
            "test",
            "tests",
            "covered",
            "coverage",
            "behavior",
            "path",
            "case",
            "cases",
        }
        for criterion in acceptance_criteria:
            tokens = [
                token
                for token in re.findall(r"[a-zA-Z][a-zA-Z0-9_]{3,}", criterion.lower())
                if token not in generic_tokens
            ]
            if tokens and not any(token in joined_tests for token in tokens):
                missing.append(criterion)
        return missing

    async def _active_rules(self) -> list[str]:
        conn = await self._db.connect()
        cursor = await conn.execute(
            """
            SELECT rule_text
            FROM rules
            WHERE enabled = 1
            ORDER BY priority DESC, updated_at DESC
            LIMIT 20
            """
        )
        rows = await cursor.fetchall()
        return [str(row["rule_text"]) for row in rows]

    def _build_context_pack(
        self,
        *,
        title: str,
        description: str,
        acceptance_criteria: list[str],
        evidence_entries: list[EvidenceBoardEntry],
        impacted_files: list[str],
        related_tests: list[str],
        missing_coverage: list[str],
        active_rules: list[str],
    ) -> str:
        evidence_lines = [
            "- "
            f"{entry.source_type} (trust={entry.trust_level}) "
            f"[{entry.extraction_status}] — "
            f"{entry.source_path or entry.source_url or 'n/a'}"
            for entry in evidence_entries
        ] or ["- none"]
        findings_lines = [
            f"- [{entry.source_type}:{entry.source_path or entry.source_url or 'n/a'}] {finding}"
            for entry in evidence_entries
            for finding in entry.findings
        ] or ["- none"]
        impacted_lines = [f"- {path}" for path in impacted_files] or ["- none"]
        test_lines = [f"- {path}" for path in related_tests] or ["- none"]
        missing_lines = [f"- {line}" for line in missing_coverage] or ["- none"]
        rule_lines = [f"- {line}" for line in active_rules] or ["- none"]
        acceptance_lines = [f"- {line}" for line in acceptance_criteria] or ["- none"]

        return "\n".join(
            [
                f"# Investigation: {title}",
                "",
                "## Source Work Item",
                description or "No work item description provided.",
                "",
                "### Acceptance Criteria",
                *acceptance_lines,
                "",
                "## Evidence Sources",
                *evidence_lines,
                "",
                "## Extracted Findings",
                *findings_lines,
                "",
                "## Impacted Code Areas",
                *impacted_lines,
                "",
                "## Related Tests",
                *test_lines,
                "",
                "## Missing Test Coverage",
                *missing_lines,
                "",
                "## Active Rules",
                *rule_lines,
                "",
                "## Constraints",
                "- Do not modify unrelated modules.",
                "- Do not auto-apply generated patches.",
                "",
                "## Expected Output",
                "- Root-cause analysis.",
                "- Implementation plan.",
                "- Test plan.",
            ]
        )

    def _write_context_pack(self, context_pack: str) -> str:
        context_pack_dir = self._config.memopilot_dir / "context-packs"
        context_pack_dir.mkdir(parents=True, exist_ok=True)
        file_name = f"investigation-{datetime.now(UTC).strftime('%Y%m%d-%H%M%S')}.md"
        path = context_pack_dir / file_name
        path.write_text(context_pack, encoding="utf-8")
        return str(path)

    async def _summarize_findings(
        self,
        *,
        conn: aiosqlite.Connection,
        investigation_session_id: str,
        workspace_root: str,
    ) -> InvestigationFindingsSummary:
        await self._require_session(conn, session_id=investigation_session_id)
        evidence_entries = await self.list_evidence_board(
            task_run_id=None,
            investigation_session_id=investigation_session_id,
            workspace_root=workspace_root,
            conn=conn,
        )
        findings = self._dedupe_preserve_order(
            finding.strip()
            for entry in evidence_entries
            for finding in entry.findings
            if finding and finding.strip()
        )
        if not findings:
            raise ValueError("No evidence findings available for investigation session")

        extracted_files: list[str] = []
        for finding in findings:
            extracted_files.extend(self._extract_code_paths(finding))
        impacted_files = self._dedupe_preserve_order(
            extracted_files + await self._discover_impacted_files(findings)
        )[:10]
        impacted_callers = self._extract_impacted_callers(findings)
        root_cause = self._extract_root_cause_findings(findings)
        acceptance_criteria = self._extract_acceptance_criteria(
            findings,
            impacted_files=impacted_files,
            impacted_callers=impacted_callers,
            root_cause=root_cause,
        )
        task_run_id = await self._resolve_investigation_task_run_id(
            conn=conn,
            investigation_session_id=investigation_session_id,
        )
        return InvestigationFindingsSummary(
            findings=findings,
            root_cause=root_cause,
            impacted_files=impacted_files,
            impacted_callers=impacted_callers,
            acceptance_criteria=acceptance_criteria,
            task_run_id=task_run_id,
        )

    async def _resolve_investigation_task_run_id(
        self,
        *,
        conn: aiosqlite.Connection,
        investigation_session_id: str,
    ) -> str | None:
        cursor = await conn.execute(
            """
            SELECT task_run_id
            FROM evidence_sources
            WHERE investigation_session_id = ? AND task_run_id IS NOT NULL
            GROUP BY task_run_id
            ORDER BY MAX(created_at) DESC
            LIMIT 2
            """,
            (investigation_session_id,),
        )
        rows = await cursor.fetchall()
        task_run_ids = [str(row["task_run_id"]) for row in rows if row["task_run_id"]]
        if len(task_run_ids) == 1:
            return task_run_ids[0]

        cursor = await conn.execute(
            """
            SELECT id
            FROM task_runs
            WHERE investigation_session_id = ?
            ORDER BY created_at DESC
            LIMIT 2
            """,
            (investigation_session_id,),
        )
        rows = await cursor.fetchall()
        session_task_runs = [str(row["id"]) for row in rows if row["id"]]
        if len(session_task_runs) == 1:
            return session_task_runs[0]
        return None

    def _build_plan_steps(self, summary: InvestigationFindingsSummary) -> list[PlanStep]:
        steps: list[PlanStep] = []
        seen: set[tuple[str, str | None, str | None]] = set()

        def append_step(
            description: str,
            *,
            target_file: str | None = None,
            target_symbol: str | None = None,
        ) -> None:
            clean_description = description.strip()
            if not clean_description:
                return
            key = (clean_description.lower(), target_file, target_symbol)
            if key in seen:
                return
            seen.add(key)
            steps.append(
                PlanStep(
                    step_number=len(steps) + 1,
                    description=clean_description,
                    target_file=target_file,
                    target_symbol=target_symbol,
                )
            )

        primary_file = summary.impacted_files[0] if summary.impacted_files else None
        if summary.root_cause:
            append_step(
                f"Fix the root cause: {self._truncate_text(summary.root_cause[0])}",
                target_file=primary_file,
            )

        for path in summary.impacted_files[:3]:
            append_step(
                f"Update {path} to reflect the investigation findings",
                target_file=path,
            )

        if summary.impacted_callers:
            append_step(
                "Verify impacted callers remain correct: "
                f"{self._truncate_text(summary.impacted_callers[0])}",
                target_file=self._extract_first_code_path(summary.impacted_callers)
                or primary_file,
            )

        for criterion in summary.acceptance_criteria[:2]:
            append_step(f"Add or update validation for: {self._truncate_text(criterion)}")

        if not steps:
            append_step("Review the investigation findings and implement the smallest safe fix")
        return steps

    def _build_plan_task_description(
        self,
        session_row: aiosqlite.Row,
        summary: InvestigationFindingsSummary,
    ) -> str:
        parts = [f"Investigation: {session_row['title']}"]
        description = str(session_row["description"] or "").strip()
        if description:
            parts.append(f"Description: {description}")
        if summary.root_cause:
            parts.append("Root cause findings: " + "; ".join(summary.root_cause[:2]))
        if summary.acceptance_criteria:
            parts.append("Acceptance criteria: " + "; ".join(summary.acceptance_criteria[:3]))
        return "\n".join(parts)

    async def _store_confirmed_memory_item(
        self,
        *,
        conn: aiosqlite.Connection,
        investigation_session_id: str,
        workspace_root: str,
        finding_kind: str,
        title: str,
        body: str,
        memory_class: str,
        trust_level: int,
    ) -> str:
        cursor = await conn.execute(
            """
            SELECT id
            FROM memory_items
            WHERE source = 'investigation'
              AND json_extract(tags_json, '$.investigation_session_id') = ?
              AND json_extract(tags_json, '$.finding_kind') = ?
            LIMIT 1
            """,
            (investigation_session_id, finding_kind),
        )
        existing = await cursor.fetchone()
        now = datetime.now(UTC).isoformat()
        tags_json = json.dumps(
            {
                "approved": True,
                "pending_approval": False,
                "finding_kind": finding_kind,
                "investigation_session_id": investigation_session_id,
            }
        )

        if existing is not None:
            await conn.execute(
                """
                UPDATE memory_items
                SET title = ?,
                    body = ?,
                    trust_level = ?,
                    tags_json = ?,
                    stale = 0,
                    memory_class = ?,
                    memory_status = 'confirmed',
                    visibility_scope = 'workspace',
                    reusable = 1,
                    review_required = 0,
                    workspace_root = ?,
                    updated_at = ?
                WHERE id = ?
                """,
                (
                    title,
                    body,
                    trust_level,
                    tags_json,
                    memory_class,
                    workspace_root,
                    now,
                    existing["id"],
                ),
            )
            return str(existing["id"])

        memory_item_id = uuid.uuid4().hex
        await conn.execute(
            """
            INSERT INTO memory_items
            (
                id, type, title, body, source, source_path, source_hash,
                trust_level, tags_json, stale, memory_class, memory_status,
                visibility_scope, reusable, review_required, workspace_root,
                created_at, updated_at
            )
            VALUES (
                ?, 'ai_summary', ?, ?, 'investigation', NULL, NULL,
                ?, ?, 0, ?, 'confirmed',
                'workspace', 1, 0, ?,
                ?, ?
            )
            """,
            (
                memory_item_id,
                title,
                body,
                trust_level,
                tags_json,
                memory_class,
                workspace_root,
                now,
                now,
            ),
        )
        return memory_item_id

    def _normalize_workspace_root(self, workspace_root: str | None) -> str:
        return str(Path(workspace_root or self._config.workspace_path).resolve())

    def _extract_root_cause_findings(self, findings: list[str]) -> list[str]:
        keywords = (
            "root cause",
            "caused by",
            "because",
            "due to",
            "failure",
            "fails",
            "failing",
            "error",
            "exception",
            "timeout",
            "regression",
            "invalid",
            "missing",
            "broken",
        )
        matches = [
            finding
            for finding in findings
            if any(keyword in finding.lower() for keyword in keywords)
        ]
        return self._dedupe_preserve_order(matches or findings[:2])[:4]

    def _extract_impacted_callers(self, findings: list[str]) -> list[str]:
        matches = [
            finding
            for finding in findings
            if any(token in finding.lower() for token in ("caller", "called by", "invoked by"))
        ]
        return self._dedupe_preserve_order(matches)[:5]

    def _extract_acceptance_criteria(
        self,
        findings: list[str],
        *,
        impacted_files: list[str],
        impacted_callers: list[str],
        root_cause: list[str],
    ) -> list[str]:
        matches = [
            finding
            for finding in findings
            if any(
                token in finding.lower()
                for token in (
                    "acceptance",
                    "should",
                    "must",
                    "expected",
                    "ensure",
                    "verify",
                    "covered",
                    "test case:",
                )
            )
        ]
        criteria = self._dedupe_preserve_order(matches)[:5]
        if criteria:
            return criteria

        synthesized: list[str] = []
        if root_cause:
            synthesized.append(
                "Eliminate the investigated failure condition: "
                f"{self._truncate_text(root_cause[0], limit=140)}"
            )
        if impacted_files:
            synthesized.append(
                f"Constrain the patch to the impacted code paths: {', '.join(impacted_files[:3])}"
            )
        if impacted_callers:
            synthesized.append("Keep impacted callers compatible with the fix.")
        if not synthesized:
            synthesized.append(
                "Address the investigation findings without regressing nearby behavior."
            )
        return synthesized

    def _extract_code_paths(self, text: str) -> list[str]:
        pattern = re.compile(
            r"[A-Za-z0-9_./\\-]+\.(?:py|js|jsx|ts|tsx|java|go|rb|php|cs|cpp|c|h|hpp|sql)"
        )
        return self._dedupe_preserve_order(
            match.group(0).strip(".,:;()[]{}<>")
            for match in pattern.finditer(text)
        )

    def _extract_first_code_path(self, lines: list[str]) -> str | None:
        for line in lines:
            paths = self._extract_code_paths(line)
            if paths:
                return paths[0]
        return None

    def _dedupe_preserve_order(self, values: list[str] | tuple[str, ...] | Any) -> list[str]:
        seen: set[str] = set()
        deduped: list[str] = []
        for value in values:
            text = str(value).strip()
            if not text or text in seen:
                continue
            seen.add(text)
            deduped.append(text)
        return deduped

    def _truncate_text(self, text: str, *, limit: int = 160) -> str:
        stripped = text.strip()
        if len(stripped) <= limit:
            return stripped
        return stripped[: limit - 3].rstrip() + "..."

    def _bullet_list(self, values: list[str]) -> str:
        return "\n".join(f"- {value}" for value in values)

    def _parse_findings(self, raw: str | None) -> list[str]:
        if not raw:
            return []
        value: Any = json.loads(raw)
        if isinstance(value, dict):
            findings = value.get("findings", [])
            if isinstance(findings, list):
                return [str(item) for item in findings]
            return []
        if isinstance(value, list):
            return [str(item) for item in value]
        return []

    def _parse_extraction_status(self, source_type: str, raw: str | None) -> str:
        if not raw:
            return self._default_extraction_status(source_type)
        value: Any = json.loads(raw)
        if isinstance(value, dict):
            status = value.get("extraction_status")
            if isinstance(status, str) and status:
                return status
        return self._default_extraction_status(source_type)

    def _parse_redacted_values(self, raw: str | None) -> int:
        if not raw:
            return 0
        value: Any = json.loads(raw)
        if isinstance(value, dict):
            redacted = value.get("redacted_values")
            if isinstance(redacted, int) and redacted >= 0:
                return redacted
        return 0

    def _default_extraction_status(self, source_type: str) -> str:
        if source_type in {"image", "screenshot", "pdf_doc"}:
            return "requires_ocr"
        return "ok"

    def _is_within_workspace(self, candidate: Path, root: Path) -> bool:
        try:
            candidate.relative_to(root.resolve())
            return True
        except ValueError:
            return False

    def _contains_parent_reference(self, candidate: Path) -> bool:
        return any(part == ".." for part in candidate.parts)

    def _important_tokens(self, findings: list[str]) -> list[str]:
        token_set: set[str] = set()
        for finding in findings:
            for token in re.findall(r"[a-zA-Z_][a-zA-Z0-9_]{3,}", finding.lower()):
                token_set.add(token)
        return sorted(token_set)[:80]
