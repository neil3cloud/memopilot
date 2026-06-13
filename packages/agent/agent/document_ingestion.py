"""Document ingestion for PDF, Excel, and CSV evidence."""

from __future__ import annotations

import csv
import hashlib
import io
import logging
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)


@dataclass
class DocumentChunk:
    """A chunk of extracted document content."""

    chunk_index: int
    chunk_text: str
    source_hash: str = ""
    trust_level: int = 3
    memory_class: str = "evidence"
    memory_status: str = "evidence_only"


@dataclass
class ExtractionResult:
    """Result of document extraction."""

    source_type: str
    chunks: list[DocumentChunk] = field(default_factory=list)
    metadata: dict = field(default_factory=dict)
    error: str | None = None
    requires_ocr: bool = False


def extract_pdf(file_path: Path, chunk_size: int = 2000) -> ExtractionResult:
    """Extract text from a PDF file.

    - Text-based PDFs: extract per page, chunk at chunk_size chars
    - Scanned PDFs: detect and flag, don't attempt OCR in v1.5
    """

    try:
        import pdfplumber
    except ImportError:
        return ExtractionResult(source_type="pdf_doc", error="pdfplumber not installed")

    chunks: list[DocumentChunk] = []
    total_text = ""
    tables_found = 0
    page_count = 0
    source_hash = _hash_file(file_path)

    try:
        with pdfplumber.open(str(file_path)) as pdf:
            page_count = len(pdf.pages)
            for page_num, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    total_text += text + "\n"

                tables = page.extract_tables()
                tables_found += len(tables)
                for table in tables:
                    table_text = "\n".join(
                        " | ".join(str(cell or "") for cell in row) for row in table if row
                    )
                    if table_text.strip():
                        total_text += f"\n[Table from page {page_num}]\n{table_text}\n"

        if page_count > 0 and len(total_text.strip()) < 50 * page_count:
            return ExtractionResult(
                source_type="pdf_doc",
                requires_ocr=True,
                metadata={"pages": page_count, "tables": tables_found},
                error="Scanned PDF detected. OCR support is planned for v2.",
            )

        chunks = _chunk_text(total_text, chunk_size=chunk_size, source_hash=source_hash)
        return ExtractionResult(
            source_type="pdf_doc",
            chunks=chunks,
            metadata={"pages": page_count, "tables": tables_found},
        )
    except Exception as exc:
        logger.warning("PDF extraction failed for %s: %s", file_path, exc)
        return ExtractionResult(source_type="pdf_doc", error=str(exc))


def extract_excel(
    file_path: Path,
    sheet_names: list[str] | None = None,
    column_mapping: dict[str, str] | None = None,
) -> ExtractionResult:
    """Extract data from an Excel file.

    - Lists sheets if sheet_names is None
    - Extracts selected sheets with header detection
    - Applies column mapping if provided
    """

    try:
        import openpyxl
    except ImportError:
        return ExtractionResult(source_type="spreadsheet", error="openpyxl not installed")

    source_hash = _hash_file(file_path)
    try:
        wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)

        if sheet_names is None:
            available = list(wb.sheetnames)
            wb.close()
            return ExtractionResult(
                source_type="spreadsheet",
                metadata={"available_sheets": available},
            )

        chunks: list[DocumentChunk] = []
        for sheet_name in sheet_names:
            if sheet_name not in wb.sheetnames:
                continue
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue

            headers = [str(h or f"Column_{i}") for i, h in enumerate(rows[0], start=1)]
            if column_mapping:
                headers = [column_mapping.get(h, h) for h in headers]

            for row_idx, row in enumerate(rows[1:], start=1):
                row_data = {
                    headers[i]: str(cell or "") for i, cell in enumerate(row) if i < len(headers)
                }
                row_text = " | ".join(f"{key}: {value}" for key, value in row_data.items() if value)
                if row_text.strip():
                    chunks.append(
                        DocumentChunk(
                            chunk_index=len(chunks),
                            chunk_text=f"[{sheet_name} row {row_idx}] {row_text}",
                            source_hash=source_hash,
                            trust_level=3,
                        )
                    )

        wb.close()
        return ExtractionResult(source_type="spreadsheet", chunks=chunks)
    except Exception as exc:
        logger.warning("Excel extraction failed for %s: %s", file_path, exc)
        return ExtractionResult(source_type="spreadsheet", error=str(exc))


def extract_docx(file_path: Path, chunk_size: int = 2000) -> ExtractionResult:
    """Extract text from a Word document by heading sections and tables."""
    try:
        from docx import Document
    except ImportError:
        return ExtractionResult(source_type="word_doc", error="python-docx not installed")

    source_hash = _hash_file(file_path)
    section_chunks: list[str] = []
    table_count = 0
    paragraph_count = 0
    current_section: list[str] = []
    current_heading = "Document"

    def flush_section() -> None:
        if current_section:
            section_chunks.append(f"[{current_heading}]\n" + "\n".join(current_section))

    try:
        document = Document(str(file_path))
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            paragraph_count += 1
            style_name = (getattr(paragraph.style, "name", "") or "").lower()
            if style_name.startswith("heading"):
                flush_section()
                current_heading = text
                current_section = []
                continue
            current_section.append(text)

        flush_section()

        for table_index, table in enumerate(document.tables, start=1):
            rows: list[str] = []
            for row in table.rows:
                values = [cell.text.strip() for cell in row.cells]
                row_text = " | ".join(value for value in values if value)
                if row_text:
                    rows.append(row_text)
            if rows:
                table_count += 1
                section_chunks.append(f"[Table {table_index}]\n" + "\n".join(rows))

        combined_text = "\n\n".join(section_chunks)
        chunks = _chunk_text(combined_text, chunk_size=chunk_size, source_hash=source_hash)
        return ExtractionResult(
            source_type="word_doc",
            chunks=chunks,
            metadata={
                "paragraphs": paragraph_count,
                "tables": table_count,
                "sections": len(section_chunks),
            },
        )
    except Exception as exc:
        logger.warning("Word extraction failed for %s: %s", file_path, exc)
        return ExtractionResult(source_type="word_doc", error=str(exc))


def extract_pptx(file_path: Path, chunk_size: int = 2000) -> ExtractionResult:
    """Extract text from PowerPoint by slide and flag skipped embedded images."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return ExtractionResult(source_type="powerpoint_doc", error="python-pptx not installed")

    source_hash = _hash_file(file_path)
    slide_chunks: list[str] = []
    skipped_images = 0

    try:
        presentation = Presentation(str(file_path))
        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines: list[str] = []
            slide_image_count = 0
            for shape in slide.shapes:
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                    slide_image_count += 1
                    skipped_images += 1
                    continue
                text = getattr(shape, "text", "")
                if text and text.strip():
                    lines.append(text.strip())
            if slide_image_count:
                plural = "s" if slide_image_count != 1 else ""
                lines.append(f"[Embedded image{plural} skipped: {slide_image_count}]")
            if lines:
                slide_chunks.append(f"[Slide {slide_index}]\n" + "\n".join(lines))

        combined_text = "\n\n".join(slide_chunks)
        chunks = _chunk_text(combined_text, chunk_size=chunk_size, source_hash=source_hash)
        return ExtractionResult(
            source_type="powerpoint_doc",
            chunks=chunks,
            metadata={"slides": len(presentation.slides), "skipped_images": skipped_images},
        )
    except Exception as exc:
        logger.warning("PowerPoint extraction failed for %s: %s", file_path, exc)
        return ExtractionResult(source_type="powerpoint_doc", error=str(exc))


def detect_csv_delimiter(content: str) -> str:
    """Auto-detect CSV delimiter."""

    try:
        dialect = csv.Sniffer().sniff(content[:4096])
        return dialect.delimiter
    except csv.Error:
        return ","


def extract_csv(
    file_path: Path,
    delimiter: str | None = None,
    column_mapping: dict[str, str] | None = None,
) -> ExtractionResult:
    """Extract data from a CSV file with auto-delimiter detection."""

    try:
        content = file_path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        logger.warning("CSV extraction failed for %s: %s", file_path, exc)
        return ExtractionResult(source_type="csv_data", error=str(exc))

    if delimiter is None:
        delimiter = detect_csv_delimiter(content)

    chunks: list[DocumentChunk] = []
    source_hash = _hash_file(file_path)
    reader = csv.DictReader(io.StringIO(content), delimiter=delimiter)

    for row_idx, row in enumerate(reader):
        normalized_row = dict(row)
        if column_mapping:
            normalized_row = {
                column_mapping.get(key, key): value for key, value in normalized_row.items()
            }
        row_text = " | ".join(f"{key}: {value}" for key, value in normalized_row.items() if value)
        if row_text.strip():
            chunks.append(
                DocumentChunk(
                    chunk_index=len(chunks),
                    chunk_text=f"[Row {row_idx + 1}] {row_text}",
                    source_hash=source_hash,
                    trust_level=3,
                )
            )

    return ExtractionResult(
        source_type="csv_data",
        chunks=chunks,
        metadata={"delimiter": delimiter},
    )


def _chunk_text(text: str, *, chunk_size: int, source_hash: str) -> list[DocumentChunk]:
    chunks: list[DocumentChunk] = []
    for start in range(0, len(text), chunk_size):
        chunk_text = text[start : start + chunk_size]
        if chunk_text.strip():
            chunks.append(
                DocumentChunk(
                    chunk_index=len(chunks),
                    chunk_text=chunk_text,
                    source_hash=source_hash,
                    trust_level=3,
                )
            )
    return chunks


def _hash_file(file_path: Path) -> str:
    try:
        digest = hashlib.sha256(file_path.read_bytes()).hexdigest()
    except OSError:
        digest = ""
    return digest
