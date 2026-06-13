"""Tests for Group 5 investigation mode and evidence workflows."""

from __future__ import annotations

import openpyxl
import pytest
from docx import Document
from httpx import AsyncClient
from PIL import Image
from pptx import Presentation

from agent.investigation_service import InvestigationService


@pytest.mark.asyncio
async def test_evidence_source_classification_and_trust_levels(
    client: AsyncClient,
    test_token: str,
    tmp_workspace,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    fixtures = [
        ("notes.md", "# title\nInvestigation note", "markdown_doc", 4),
        (
            "trace.log",
            "Traceback (most recent call last):\n  File \"main.py\", line 1, in <module>\nValueError: boom",
            "stack_trace",
            4,
        ),
        ("incident.txt", "timeout error without stack frames", "text_log", 3),
        ("metrics.csv", "name,value\nlatency,12", "csv_data", 3),
        ("payload.json", '{"error":"invalid_token"}', "api_payload", 3),
        ("schema.sql", "CREATE TABLE users(id INTEGER);", "database_schema", 4),
        ("service.py", "def run():\n    return None\n", "existing_code", 5),
        ("specs.xlsx", None, "spreadsheet", 3),
        ("requirements.docx", None, "word_doc", 3),
        ("slides.pptx", None, "powerpoint_doc", 3),
        ("screenshot.png", None, "screenshot", 4),
        ("diagram.png", None, "image", 5),
    ]
    for name, content, expected_type, expected_trust in fixtures:
        file_path = tmp_workspace / name
        if name.endswith(".xlsx"):
            workbook = openpyxl.Workbook()
            sheet = workbook.active
            sheet.append(["title", "steps", "expected"])
            sheet.append(["Login", "Open app", "Dashboard"])
            workbook.save(file_path)
            workbook.close()
        elif name.endswith(".docx"):
            document = Document()
            document.add_paragraph("Wave 4 doc ingestion")
            document.save(file_path)
        elif name.endswith(".pptx"):
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Wave 4 slides"
            slide.placeholders[1].text = "PowerPoint text extraction"
            presentation.save(file_path)
        elif name.endswith(".png"):
            image = Image.new("RGB", (16, 16), color=(255, 255, 255))
            image.save(file_path)
        else:
            file_path.write_text(content, encoding="utf-8")

        attached = await client.post(
            "/v1/investigation/evidence/attach",
            headers=headers,
            json={"evidence_path": str(file_path)},
        )
        assert attached.status_code == 200
        body = attached.json()
        assert body["source_type"] == expected_type
        assert body["trust_level"] == expected_trust
        if expected_type in {"image", "screenshot"}:
            assert body["extraction_status"] in {"metadata_only", "ok"}
            assert any("Image analysis:" in finding for finding in body["findings"])
            if expected_type == "screenshot":
                assert any("Screenshot heuristic:" in finding for finding in body["findings"])


@pytest.mark.asyncio
async def test_attach_evidence_redacts_secrets(
    client: AsyncClient,
    test_token: str,
    tmp_workspace,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    evidence = tmp_workspace / "incident.txt"
    evidence.write_text("api_key=super-secret-value", encoding="utf-8")
    attached = await client.post(
        "/v1/investigation/evidence/attach",
        headers=headers,
        json={"evidence_path": str(evidence)},
    )
    assert attached.status_code == 200
    body = attached.json()
    assert body["redacted_values"] >= 1
    assert "super-secret-value" not in "\n".join(body["findings"])

    board = await client.get("/v1/investigation/evidence", headers=headers)
    assert board.status_code == 200
    item = board.json()["items"][0]
    assert item["extraction_status"] == "ok"
    assert item["redacted_values"] >= 1


@pytest.mark.asyncio
async def test_attach_evidence_rejects_parent_traversal(
    client: AsyncClient,
    test_token: str,
    tmp_workspace,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    outside_evidence = tmp_workspace.parent / "outside.txt"
    outside_evidence.write_text("outside workspace", encoding="utf-8")
    attached = await client.post(
        "/v1/investigation/evidence/attach",
        headers=headers,
        json={"evidence_path": "../outside.txt"},
    )
    assert attached.status_code == 400
    assert "must not traverse parent directories" in attached.json()["detail"]


@pytest.mark.asyncio
async def test_run_investigation_builds_context_pack_and_coverage(
    client: AsyncClient,
    test_token: str,
    tmp_workspace,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    (tmp_workspace / "service.py").write_text(
        "def retry_service() -> None:\n    pass\n",
        encoding="utf-8",
    )
    (tmp_workspace / "tests").mkdir()
    (tmp_workspace / "tests" / "test_service.py").write_text(
        "def test_retry_service() -> None:\n    assert True\n",
        encoding="utf-8",
    )
    await client.post("/v1/workspace/index", headers=headers)

    evidence = tmp_workspace / "summary.txt"
    evidence.write_text(
        "retry_service failed under timeout and service error path.",
        encoding="utf-8",
    )
    attached = await client.post(
        "/v1/investigation/evidence/attach",
        headers=headers,
        json={"evidence_path": str(evidence)},
    )
    assert attached.status_code == 200

    run = await client.post(
        "/v1/investigation/run",
        headers=headers,
        json={
            "title": "Retry failure",
            "description": "Investigate retry behavior regression",
            "acceptance_criteria": [
                "service retry path is covered by tests",
                "timeout fallback behavior is covered by tests",
            ],
        },
    )
    assert run.status_code == 200
    payload = run.json()
    assert payload["evidence_count"] >= 1
    assert any(path.endswith("service.py") for path in payload["impacted_files"])
    assert any("test_service.py" in path for path in payload["related_tests"])
    assert "timeout fallback behavior is covered by tests" in payload["missing_test_coverage"]
    assert "## Evidence Sources" in payload["context_pack"]
    assert "## Missing Test Coverage" in payload["context_pack"]
    assert "[text_log:" in payload["context_pack"]


@pytest.mark.asyncio
async def test_excel_preview_and_mapping_extracts_test_cases(
    client: AsyncClient,
    test_token: str,
    tmp_workspace,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    workbook_path = tmp_workspace / "test-cases.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.append(["Case Title", "Procedure", "Expected Result", "Priority"])
    sheet.append(["Login success", "Use valid creds", "Home page opens", "P1"])
    sheet.append(["Login fail", "Use bad creds", "Error appears", "P2"])
    workbook.save(workbook_path)
    workbook.close()

    preview = await client.post(
        "/v1/investigation/evidence/columns",
        headers=headers,
        json={"evidence_path": str(workbook_path)},
    )
    assert preview.status_code == 200
    preview_payload = preview.json()
    assert preview_payload["source_type"] == "spreadsheet"
    assert "Case Title" in preview_payload["columns"]
    assert preview_payload["requires_confirmation"] is True

    attached = await client.post(
        "/v1/investigation/evidence/attach",
        headers=headers,
        json={
            "evidence_path": str(workbook_path),
            "column_mapping": {
                "title": "Case Title",
                "steps": "Procedure",
                "expected": "Expected Result",
            },
        },
    )
    assert attached.status_code == 200
    body = attached.json()
    assert body["source_type"] == "spreadsheet"
    assert any("Test case:" in finding for finding in body["findings"])
    assert any("Data dictionary:" in finding for finding in body["findings"])


@pytest.mark.asyncio
async def test_scanned_pdf_forces_trust_level_four(
    client: AsyncClient,
    test_token: str,
    tmp_workspace,
    monkeypatch: pytest.MonkeyPatch,
):
    monkeypatch.setattr(
        InvestigationService,
        "_extract_pdf_findings",
        lambda self, _path: (
            ["Scanned PDF detected. OCR is required before reliable extraction."],
            "requires_ocr",
            True,
        ),
    )

    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    evidence = tmp_workspace / "scan.pdf"
    evidence.write_text("scanned content placeholder", encoding="utf-8")
    attached = await client.post(
        "/v1/investigation/evidence/attach",
        headers=headers,
        json={"evidence_path": str(evidence)},
    )
    assert attached.status_code == 200
    payload = attached.json()
    assert payload["source_type"] == "pdf_doc"
    assert payload["extraction_status"] == "requires_ocr"
    assert payload["trust_level"] == 4


@pytest.mark.asyncio
async def test_image_analysis_extracts_metadata(
    client: AsyncClient,
    test_token: str,
    tmp_workspace,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    image_path = tmp_workspace / "error-dialog.png"
    image = Image.new("RGB", (800, 600), color=(250, 250, 250))
    image.save(image_path)

    attached = await client.post(
        "/v1/investigation/evidence/attach",
        headers=headers,
        json={"evidence_path": str(image_path)},
    )
    assert attached.status_code == 200
    body = attached.json()
    assert body["source_type"] == "screenshot"
    assert body["trust_level"] == 4
    assert body["extraction_status"] in {"metadata_only", "ok"}
    assert any("Image analysis: 800x600" in finding for finding in body["findings"])


@pytest.mark.asyncio
async def test_investigation_session_endpoints_manage_evidence(
    client: AsyncClient,
    test_token: str,
    tmp_workspace,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    evidence = tmp_workspace / "investigation-notes.txt"
    evidence.write_text("login retry timeout observed", encoding="utf-8")

    started = await client.post(
        "/v1/investigation/start",
        headers=headers,
        json={"title": "Retry timeout", "description": "Investigate login retries"},
    )
    assert started.status_code == 200
    session = started.json()
    session_id = session["id"]
    assert session["status"] == "open"
    assert session["evidence_count"] == 0

    fetched = await client.get(f"/v1/investigation/{session_id}", headers=headers)
    assert fetched.status_code == 200
    assert fetched.json()["id"] == session_id

    attached = await client.post(
        f"/v1/investigation/{session_id}/evidence",
        headers=headers,
        json={"evidence_path": str(evidence)},
    )
    assert attached.status_code == 200
    attached_body = attached.json()
    evidence_id = attached_body["evidence_id"]
    assert attached_body["investigation_session_id"] == session_id

    board = await client.get(
        "/v1/investigation/evidence",
        headers=headers,
        params={"investigation_session_id": session_id},
    )
    assert board.status_code == 200
    assert board.json()["items"][0]["evidence_id"] == evidence_id

    with_evidence = await client.get(f"/v1/investigation/{session_id}", headers=headers)
    assert with_evidence.status_code == 200
    with_evidence_body = with_evidence.json()
    assert with_evidence_body["evidence_count"] == 1
    assert with_evidence_body["evidence"][0]["evidence_id"] == evidence_id

    transitioned = await client.post(
        f"/v1/investigation/{session_id}/transition-to-patch",
        headers=headers,
    )
    assert transitioned.status_code == 200
    assert transitioned.json()["status"] == "patch_generated"

    deleted = await client.delete(
        f"/v1/investigation/{session_id}/evidence/{evidence_id}",
        headers=headers,
    )
    assert deleted.status_code == 200
    assert deleted.json() == {"evidence_id": evidence_id, "removed": True}

    after_delete = await client.get(f"/v1/investigation/{session_id}", headers=headers)
    assert after_delete.status_code == 200
    assert after_delete.json()["evidence_count"] == 0
