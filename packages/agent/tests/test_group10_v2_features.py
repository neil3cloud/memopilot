"""Tests for v2 feature additions."""

from __future__ import annotations

from pathlib import Path

import pytest
import yaml
from docx import Document
from httpx import AsyncClient
from PIL import Image
from pptx import Presentation

from agent.flow_builder import validate_flow
from agent.image_analysis import ImageAnalysisResult
from agent.workspace_roots import WorkspaceRootsService


@pytest.mark.asyncio
async def test_docx_pptx_and_image_endpoints(
    client: AsyncClient,
    test_token: str,
    tmp_workspace: Path,
    monkeypatch: pytest.MonkeyPatch,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    docx_path = tmp_workspace / "brief.docx"
    doc = Document()
    doc.add_heading("Summary", level=1)
    doc.add_paragraph("Ship the workspace isolation update.")
    doc.save(docx_path)

    pptx_path = tmp_workspace / "briefing.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Launch"
    slide.placeholders[1].text = "Enable review memory mode"
    presentation.save(pptx_path)

    image_path = tmp_workspace / "screen.png"
    Image.new("RGB", (32, 32), color=(255, 255, 255)).save(image_path)

    async def fake_analyze_image(
        _path: Path,
        allow_cloud: bool = False,
        workspace_root: str | None = None,
    ) -> ImageAnalysisResult:
        return ImageAnalysisResult(
            description="Settings dialog with save button",
            ui_elements=["Save button"],
            error_messages=[],
            ocr_text="Settings",
            source="local",
        )

    monkeypatch.setattr("agent.api.analyze_image", fake_analyze_image)

    docx_response = await client.post(
        "/v1/evidence/extract-docx",
        headers=headers,
        json={"file_path": str(docx_path)},
    )
    assert docx_response.status_code == 200
    assert any("workspace isolation update" in chunk["chunk_text"].lower() for chunk in docx_response.json()["chunks"])

    pptx_response = await client.post(
        "/v1/evidence/extract-pptx",
        headers=headers,
        json={"file_path": str(pptx_path)},
    )
    assert pptx_response.status_code == 200
    assert any("review memory mode" in chunk["chunk_text"].lower() for chunk in pptx_response.json()["chunks"])

    image_response = await client.post(
        "/v1/evidence/analyze-image",
        headers=headers,
        json={"file_path": str(image_path), "allow_cloud": False},
    )
    assert image_response.status_code == 200
    payload = image_response.json()
    assert payload["trust_level"] == 2
    assert payload["memory_status"] == "evidence_only"
    assert payload["description"] == "Settings dialog with save button"
    assert payload["ui_elements"] == ["Save button"]


@pytest.mark.asyncio
async def test_policy_loading_precedence_and_active_rules(
    client: AsyncClient,
    test_token: str,
    test_config,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    policy_dir = test_config.workspace_path / ".memopilot-policy"
    policy_dir.mkdir(parents=True, exist_ok=True)
    (policy_dir / "team.yaml").write_text(
        yaml.safe_dump(
            {
                "name": "team-allow-local",
                "enforcement_mode": "enforce",
                "active": True,
                "rules": ["allow_model: local-model"],
            }
        ),
        encoding="utf-8",
    )

    workspace_rules_dir = test_config.memopilot_dir / "rules"
    workspace_rules_dir.mkdir(parents=True, exist_ok=True)
    (workspace_rules_dir / "workspace.rules.yaml").write_text(
        "rules:\n  - deny_model: local-model\n",
        encoding="utf-8",
    )

    loaded = await client.post("/v1/policies/load", headers=headers)
    assert loaded.status_code == 200
    assert any(item["name"] == "team-allow-local" for item in loaded.json()["items"])

    active = await client.get("/v1/policies/active", headers=headers)
    assert active.status_code == 200
    active_items = active.json()["items"]
    assert any(item["rule"] == "allow_model: local-model" for item in active_items)
    assert any(item["source_kind"] == "workspace_rules" for item in active_items)

    evaluated = await client.post(
        "/v1/policies/evaluate",
        headers=headers,
        json={"stage": "model_call", "selected_model": "local-model"},
    )
    assert evaluated.status_code == 200
    assert evaluated.json()["allowed"] is True


@pytest.mark.asyncio
async def test_flow_yaml_validation_and_constraint_enforcement(
    client: AsyncClient,
    test_token: str,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    create_response = await client.post(
        "/v1/flows/local",
        headers=headers,
        json={
            "flow_yaml": yaml.safe_dump(
                {
                    "flow_id": "guarded-flow",
                    "name": "Guarded Flow",
                    "description": "Validate approvals before patches",
                    "steps": [
                        {
                            "id": "analyze-1",
                            "name": "Analyze task",
                            "action": "analyze_task",
                            "escalate_after_failures": 2,
                            "escalate_to_model": "gpt-4o-mini",
                        },
                        {
                            "id": "patch-1",
                            "name": "Generate patch",
                            "action": "generate_patch",
                            "requires_approval": True,
                        },
                    ],
                }
            )
        },
    )
    assert create_response.status_code == 200
    flow_id = create_response.json()["flow_id"]

    run_response = await client.post(
        "/v1/flows/local/run",
        headers=headers,
        json={
            "flow_id": flow_id,
            "task_text": "Apply the fix",
            "files_changed": ["src/app.py"],
            "selected_model": "local-model",
            "failure_count": 3,
            "constraints": ["no_file_modification_without_approval"],
        },
    )
    assert run_response.status_code == 200
    payload = run_response.json()
    assert payload["status"] == "blocked"
    assert payload["blocked_reason"] == "Approval required for step 'Generate patch'"
    assert payload["steps"][0]["model_escalated_to"] == "gpt-4o-mini"


def test_validate_flow_reports_schema_and_destructive_command_errors():
    valid, errors = validate_flow(
        yaml.safe_dump(
            {
                "flow_id": "unsafe-flow",
                "name": "Unsafe Flow",
                "steps": [
                    {"action": "analyze_task"},
                    {
                        "name": "Delete everything",
                        "action": "run_validation",
                        "command": r"Remove-Item -Recurse -Force C:\repo",
                    },
                ],
            }
        )
    )
    assert valid is False
    assert any("steps[0].name is required" in error for error in errors)
    assert any("destructive command" in error.lower() for error in errors)


@pytest.mark.asyncio
async def test_workspace_roots_service_supports_active_root_switching(
    client: AsyncClient,
    test_token: str,
    test_config,
    test_db,
    tmp_workspace: Path,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    service = WorkspaceRootsService(config=test_config, db=test_db)
    initial_active = await service.get_active_root()
    assert initial_active == str(tmp_workspace.resolve())

    workspace_b = tmp_workspace.parent / "workspace-b"
    workspace_b.mkdir(parents=True, exist_ok=True)
    await service.add_workspace_root(
        root_path=str(workspace_b),
        label="Workspace B",
        activate=False,
        workspace_root=str(tmp_workspace),
    )

    switched = await service.set_active_root(str(workspace_b), workspace_root=str(tmp_workspace))
    assert switched.active is True
    assert await service.get_active_root() == str(workspace_b.resolve())

    listed = await service.list_roots(workspace_root=str(workspace_b))
    assert any(item.root_path == str(workspace_b.resolve()) and item.active for item in listed)


@pytest.mark.asyncio
async def test_workspace_memory_isolation_and_review_memory_mode(
    client: AsyncClient,
    test_token: str,
    tmp_workspace: Path,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    workspace_b = tmp_workspace.parent / "workspace-b"
    workspace_b.mkdir(parents=True, exist_ok=True)
    added = await client.post(
        "/v1/workspaces",
        headers=headers,
        json={"root_path": str(workspace_b), "label": "Workspace B"},
    )
    assert added.status_code == 200

    first = await client.post(
        "/v1/memory/writeback",
        headers=headers,
        json={"title": "Workspace A lesson", "body": "shared keyword", "workspace_root": str(tmp_workspace)},
    )
    assert first.status_code == 200
    second = await client.post(
        "/v1/memory/writeback",
        headers=headers,
        json={"title": "Workspace B lesson", "body": "shared keyword", "workspace_root": str(workspace_b)},
    )
    assert second.status_code == 200

    memory_a = await client.get(
        "/v1/memory/items",
        headers=headers,
        params={"workspace_root": str(tmp_workspace)},
    )
    assert memory_a.status_code == 200
    titles_a = [item["title"] for item in memory_a.json()["items"]]
    assert "Workspace A lesson" in titles_a
    assert "Workspace B lesson" not in titles_a

    review_evidence = await client.post(
        "/v1/reviews/evidence",
        headers=headers,
        json={"pr_number": 42, "body": "Always assert the fallback path", "workspace_root": str(tmp_workspace)},
    )
    assert review_evidence.status_code == 200
    evidence_id = review_evidence.json()["evidence_id"]

    approved = await client.post(
        "/v1/reviews/approve-lesson",
        headers=headers,
        json={
            "evidence_id": evidence_id,
            "lesson_title": "Review lesson",
            "lesson_body": "Assert fallback behavior in future fixes.",
            "workspace_root": str(tmp_workspace),
        },
    )
    assert approved.status_code == 200
    memory_item_id = approved.json()["memory_item_id"]

    review_memory = await client.get(
        "/v1/memory/items",
        headers=headers,
        params={"workspace_root": str(tmp_workspace)},
    )
    assert review_memory.status_code == 200
    item = next(entry for entry in review_memory.json()["items"] if entry["id"] == memory_item_id)
    assert item["memory_class"] == "decision"
    assert item["memory_status"] == "confirmed"
    assert item["reusable"] is True


@pytest.mark.asyncio
async def test_review_lesson_extract_and_approve_endpoints(
    client: AsyncClient,
    test_token: str,
    tmp_workspace: Path,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    extracted = await client.post(
        "/v1/memory/review-lessons/extract",
        headers=headers,
        json={
            "review_comments": [
                {
                    "body": "Always add regression coverage for fallback behavior.",
                    "path": "src/app.py",
                    "pr_url": "https://example.test/pr/42",
                    "author": "maintainer",
                },
                {
                    "body": "nit: spacing",
                    "path": "src/app.py",
                },
            ]
        },
    )
    assert extracted.status_code == 200
    lessons = extracted.json()["lessons"]
    assert len(lessons) == 1
    assert lessons[0]["source_reviewer"] == "maintainer"

    approved = await client.post(
        "/v1/memory/review-lessons/approve",
        headers=headers,
        json={
            "summary": lessons[0]["summary"],
            "context": lessons[0]["context"],
            "source_pr": lessons[0]["source_pr"],
            "source_reviewer": lessons[0]["source_reviewer"],
            "workspace_root": str(tmp_workspace),
        },
    )
    assert approved.status_code == 200
    assert approved.json()["approved"] is True

    memory_items = await client.get(
        "/v1/memory/items",
        headers=headers,
        params={"workspace_root": str(tmp_workspace)},
    )
    assert memory_items.status_code == 200
    item = next(entry for entry in memory_items.json()["items"] if entry["id"] == approved.json()["memory_item_id"])
    assert item["title"].startswith("Review lesson:")
    assert item["memory_class"] == "decision"
    assert item["memory_status"] == "confirmed"
    assert item["reusable"] is True
