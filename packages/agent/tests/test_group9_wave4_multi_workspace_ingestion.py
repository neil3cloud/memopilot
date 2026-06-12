"""Tests for Wave 4 multi-workspace support and office ingestion."""

from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document
from httpx import AsyncClient
from pptx import Presentation


@pytest.mark.asyncio
async def test_word_and_powerpoint_evidence_ingestion(
    client: AsyncClient,
    test_token: str,
    tmp_workspace: Path,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    docx_path = tmp_workspace / "release-notes.docx"
    doc = Document()
    doc.add_paragraph("Release notes for sprint 14")
    doc.add_paragraph("Fix authentication timeout edge cases")
    doc.save(docx_path)

    pptx_path = tmp_workspace / "architecture.pptx"
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "System Architecture"
    slide.placeholders[1].text = "Introduce workspace-level isolation boundaries"
    presentation.save(pptx_path)

    attached_docx = await client.post(
        "/v1/investigation/evidence/attach",
        headers=headers,
        json={"evidence_path": str(docx_path)},
    )
    assert attached_docx.status_code == 200
    docx_body = attached_docx.json()
    assert docx_body["source_type"] == "word_doc"
    assert docx_body["extraction_status"] == "ok"
    assert any("authentication timeout" in line.lower() for line in docx_body["findings"])

    attached_pptx = await client.post(
        "/v1/investigation/evidence/attach",
        headers=headers,
        json={"evidence_path": str(pptx_path)},
    )
    assert attached_pptx.status_code == 200
    pptx_body = attached_pptx.json()
    assert pptx_body["source_type"] == "powerpoint_doc"
    assert pptx_body["extraction_status"] == "ok"
    assert any("Slide 1:" in line for line in pptx_body["findings"])


@pytest.mark.asyncio
async def test_workspace_roots_support_active_root_resolution(
    client: AsyncClient,
    test_token: str,
    tmp_workspace: Path,
):
    headers = {"X-Agent-Token": test_token}
    await client.post("/v1/workspace/init", headers=headers)

    initial = await client.get("/v1/workspaces", headers=headers)
    assert initial.status_code == 200
    initial_items = initial.json()["items"]
    assert len(initial_items) >= 1
    assert any(item["active"] for item in initial_items)

    other_workspace = tmp_workspace.parent / "workspace-b"
    other_workspace.mkdir(parents=True, exist_ok=True)
    evidence_file = other_workspace / "incident.txt"
    evidence_file.write_text("Cross-workspace evidence is now allowed.", encoding="utf-8")

    added = await client.post(
        "/v1/workspaces",
        headers=headers,
        json={"root_path": str(other_workspace), "label": "Workspace B", "activate": True},
    )
    assert added.status_code == 200
    added_body = added.json()
    assert added_body["active"] is True

    listed = await client.get("/v1/workspaces", headers=headers)
    assert listed.status_code == 200
    assert any(item["label"] == "Workspace B" for item in listed.json()["items"])

    attached_relative = await client.post(
        "/v1/investigation/evidence/attach",
        headers=headers,
        json={"evidence_path": "incident.txt"},
    )
    assert attached_relative.status_code == 200
    relative_body = attached_relative.json()
    assert relative_body["source_path"] is not None
    assert str(other_workspace) in relative_body["source_path"]
